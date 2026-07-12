"""Export service: CSV, Excel, PDF, HTML, PowerPoint."""

import os
import uuid
import logging
from typing import List, Dict, Any
from datetime import datetime
import pandas as pd

from app.core.config import settings

logger = logging.getLogger(__name__)


class ExportService:
    def export(
        self,
        rows: List[Dict],
        columns: List[str],
        format: str,
        title: str = "Query Results",
        sql: str = "",
    ) -> str:
        """Export data to file and return filename."""
        os.makedirs(settings.EXPORT_DIR, exist_ok=True)
        filename = f"export_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        df = pd.DataFrame(rows, columns=columns)

        if format == "csv":
            return self._export_csv(df, filename)
        elif format == "excel":
            return self._export_excel(df, filename, title, sql)
        elif format == "pdf":
            return self._export_pdf(df, filename, title, sql)
        elif format == "html":
            return self._export_html(df, filename, title, sql)
        elif format == "pptx":
            return self._export_pptx(df, filename, title, sql)
        else:
            raise ValueError(f"Unsupported export format: {format}")

    def _export_csv(self, df: pd.DataFrame, filename: str) -> str:
        path = os.path.join(settings.EXPORT_DIR, f"{filename}.csv")
        df.to_csv(path, index=False)
        return f"{filename}.csv"

    def _export_excel(self, df: pd.DataFrame, filename: str, title: str, sql: str) -> str:
        path = os.path.join(settings.EXPORT_DIR, f"{filename}.xlsx")
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            df.to_excel(writer, sheet_name="Results", index=False)
            wb = writer.book
            ws = writer.sheets["Results"]

            # Style header row
            from openpyxl.styles import Font, PatternFill, Alignment
            header_fill = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            for cell in ws[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")

            # Auto-size columns
            for col in ws.columns:
                max_len = max(len(str(cell.value or "")) for cell in col)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

            # Add metadata sheet
            meta_ws = wb.create_sheet("Metadata")
            meta_ws.append(["Report Title", title])
            meta_ws.append(["Generated At", datetime.now().isoformat()])
            meta_ws.append(["Row Count", len(df)])
            meta_ws.append(["Column Count", len(df.columns)])
            if sql:
                meta_ws.append(["SQL Query", sql])

        return f"{filename}.xlsx"

    def _export_pdf(self, df: pd.DataFrame, filename: str, title: str, sql: str) -> str:
        path = os.path.join(settings.EXPORT_DIR, f"{filename}.pdf")
        try:
            from reportlab.lib.pagesizes import landscape, A4
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
            from reportlab.lib.units import inch

            doc = SimpleDocTemplate(path, pagesize=landscape(A4))
            styles = getSampleStyleSheet()
            story = []

            # Title
            story.append(Paragraph(title, styles["Title"]))
            story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles["Normal"]))
            story.append(Spacer(1, 0.2 * inch))

            # Table
            preview_df = df.head(500)
            data = [list(preview_df.columns)] + preview_df.fillna("").values.tolist()
            data = [[str(cell)[:50] for cell in row] for row in data]

            t = Table(data)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E3A5F")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F5F5F5")]),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]))
            story.append(t)

            if sql:
                story.append(Spacer(1, 0.2 * inch))
                story.append(Paragraph(f"SQL: {sql[:500]}", styles["Code"]))

            doc.build(story)
        except ImportError:
            # Fallback: HTML-based PDF
            html_path = self._export_html(df, filename + "_temp", title, sql)
            full_html = os.path.join(settings.EXPORT_DIR, html_path)
            # Just rename HTML → PDF as fallback
            import shutil
            shutil.copy(full_html, path.replace(".pdf", ".html"))
            return f"{filename}.html"

        return f"{filename}.pdf"

    def _export_html(self, df: pd.DataFrame, filename: str, title: str, sql: str) -> str:
        path = os.path.join(settings.EXPORT_DIR, f"{filename}.html")
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  body {{ font-family: 'Segoe UI', sans-serif; margin: 32px; background: #f9fafb; color: #1f2937; }}
  h1 {{ color: #1e3a5f; margin-bottom: 4px; }}
  .meta {{ color: #6b7280; font-size: 14px; margin-bottom: 24px; }}
  table {{ width: 100%; border-collapse: collapse; background: white; border-radius: 8px; overflow: hidden; box-shadow: 0 1px 3px rgba(0,0,0,0.1); }}
  th {{ background: #1e3a5f; color: white; padding: 10px 14px; text-align: left; font-weight: 600; }}
  td {{ padding: 9px 14px; border-bottom: 1px solid #e5e7eb; font-size: 13px; }}
  tr:nth-child(even) td {{ background: #f8fafc; }}
  tr:hover td {{ background: #eff6ff; }}
  .sql-block {{ background: #1e293b; color: #e2e8f0; padding: 16px; border-radius: 6px; font-family: monospace; font-size: 13px; margin-top: 24px; white-space: pre-wrap; word-break: break-all; }}
</style>
</head>
<body>
<h1>{title}</h1>
<div class="meta">Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} &nbsp;|&nbsp; {len(df):,} rows &times; {len(df.columns)} columns</div>
{df.to_html(index=False, classes="", border=0, na_rep="")}
{"<div class='sql-block'>" + sql + "</div>" if sql else ""}
</body>
</html>"""
        with open(path, "w", encoding="utf-8") as f:
            f.write(html)
        return f"{filename}.html"

    def _export_pptx(self, df: pd.DataFrame, filename: str, title: str, sql: str) -> str:
        path = os.path.join(settings.EXPORT_DIR, f"{filename}.pptx")
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]

            # Title slide
            slide = prs.slides.add_slide(blank)
            tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = tx.text_frame
            tf.text = title
            tf.paragraphs[0].runs[0].font.size = Pt(40)
            tf.paragraphs[0].runs[0].font.bold = True
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            sub_tx = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
            sub_tf = sub_tx.text_frame
            sub_tf.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} | {len(df):,} rows"
            sub_tf.paragraphs[0].runs[0].font.size = Pt(18)
            sub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Data slide (preview table)
            slide2 = prs.slides.add_slide(blank)
            preview = df.head(20)
            rows_count = len(preview) + 1
            cols_count = min(len(preview.columns), 8)
            preview_cols = list(preview.columns[:cols_count])

            table = slide2.shapes.add_table(
                rows_count, cols_count,
                Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.0)
            ).table

            # Header
            for ci, col in enumerate(preview_cols):
                cell = table.cell(0, ci)
                cell.text = str(col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

            # Data rows
            for ri, (_, row) in enumerate(preview.iterrows()):
                for ci, col in enumerate(preview_cols):
                    cell = table.cell(ri + 1, ci)
                    cell.text = str(row[col] if row[col] is not None else "")
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
                    if ri % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

            # Summary slide
            slide3 = prs.slides.add_slide(blank)
            tx3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tx3.text_frame.text = "Summary Statistics"
            tx3.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
            tx3.text_frame.paragraphs[0].runs[0].font.bold = True
            tx3.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

            numeric_cols = df.select_dtypes(include="number").columns.tolist()
            if numeric_cols:
                stats = df[numeric_cols[:5]].describe().round(2)
                stats_tx = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                stats_tf = stats_tx.text_frame
                stats_tf.word_wrap = True
                stats_tf.text = stats.to_string()
                stats_tf.paragraphs[0].runs[0].font.size = Pt(11)
                stats_tf.paragraphs[0].runs[0].font.name = "Courier New"

            prs.save(path)
        except ImportError:
            logger.error("python-pptx not installed")
            return self._export_html(df, filename, title, sql).replace(".html", ".html")

        return f"{filename}.pptx"
